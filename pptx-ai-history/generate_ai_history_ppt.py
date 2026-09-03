# -*- coding: utf-8 -*-
"""
生成主题为「人工智能发展史」的 PPT（5 页，每页含标题与要点）。

依赖：python-pptx (>= 0.6)
用法：python generate_ai_history_ppt.py
输出：ai_history.pptx（保存在当前工作目录）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 每页标题与要点的数据
# 结构：{ "title": 标题, "subtitle": 可选副标题/导语, "points": [要点...] }
SLIDES = [
    {
        "title": "人工智能发展史",
        "subtitle": "从图灵的思想到深度学习的浪潮",
        "points": [
            "AI（Artificial Intelligence）旨在让机器模拟人类的智能行为",
            "发展历经思想萌芽、黄金时代、两次寒冬与再度崛起",
            "本讲按时间线梳理关键里程碑、人物与技术突破",
        ],
    },
    {
        "title": "思想的萌芽与奠基（1943–1955）",
        "subtitle": "逻辑与计算为智能机器提供理论基础",
        "points": [
            "1943 年：McCulloch 与 Pitts 提出神经网络最早数学模型",
            "1950 年：图灵发表《计算机器与智能》，提出“图灵测试”",
            "同时期人工智能正式诞生前，逻辑与符号处理思想逐步成型",
        ],
    },
    {
        "title": "黄金时代与两次寒冬（1956–1993）",
        "subtitle": "从达特茅斯会议到符号主义的起伏",
        "points": [
            "1956 年：达特茅斯会议正式提出“人工智能”，标志学科诞生",
            "符号主义/专家系统取得早期成功（如 MYCIN 医疗系统）",
            "因算力与数据不足，1974–1980、1987–1993 两次进入“AI寒冬”",
        ],
    },
    {
        "title": "机器学习的复兴与深度学习崛起（1990s–2010s）",
        "subtitle": "统计方法回归，算力与数据推动新一轮突破",
        "points": [
            "支持向量机、随机森林等统计学习模型兴起",
            "2012 年：AlexNet 在 ImageNet 大胜，引爆深度学习热潮",
            "GPU、大数据与开源框架（TensorFlow 等）加速技术落地",
        ],
    },
    {
        "title": "当代 AI 浪潮与未来展望",
        "subtitle": "大模型、生成式 AI 与通用人工智能（AGI）",
        "points": [
            "Transformer 与大语言模型让 AI 进入“生成式智能”时代",
            "AlphaGo、ChatGPT 等标志性成果进入大众视野",
            "未来方向：AGI、多模态、具身智能，以及安全与伦理治理",
        ],
    },
]


def add_title_slide(prs, slide):
    """添加一张标题+要点的标准页面。"""
    layout = prs.slide_layouts[1]  # Title and Content 布局
    s = prs.slides.add_slide(layout)

    title_ph = s.shapes.title
    title_ph.text = slide["title"]

    body = s.placeholders[1]
    tf = body.text_frame
    tf.clear()

    first_used = False

    # 副标题 / 导语行 —— 占用正文首段
    if slide.get("subtitle"):
        p = tf.paragraphs[0]
        p.text = slide["subtitle"]
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            run.font.bold = True
        first_used = True

    # 逐条添加要点
    for point in slide["points"]:
        if not first_used:
            # 没有副标题时，第一条要点使用首段
            p = tf.paragraphs[0]
            first_used = True
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(18)


def main():
    prs = Presentation()
    # 设置为 16:9 宽屏，提升观感
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide in SLIDES:
        add_title_slide(prs, slide)

    out_file = "ai_history.pptx"
    prs.save(out_file)
    print("已生成：", out_file, "共", len(prs.slides._sldIdLst), "页")


if __name__ == "__main__":
    main()
